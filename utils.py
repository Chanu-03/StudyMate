# -------------------------
# utils.py — Updated for caching & faster processing
# -------------------------
import os
import json
import tempfile
import pickle
from typing import List, Tuple
from sentence_transformers import SentenceTransformer

# File libs
import docx
from pptx import Presentation
import fitz  # PyMuPDF

# Embedding model
_embedder = SentenceTransformer("all-MiniLM-L6-v2")

# -------------------------
# File text extraction
# -------------------------
def extract_text_from_pdf_bytes(file_bytes: bytes) -> str:
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = fitz.open(tmp_path)
        text = [page.get_text("text") for page in doc if page.get_text("text")]
        return "\n".join(text)
    finally:
        try:
            os.unlink(tmp_path)
        except:
            pass

def extract_text_from_docx(file_bytes: bytes) -> str:
    from io import BytesIO
    doc = docx.Document(BytesIO(file_bytes))
    return "\n".join([p.text for p in doc.paragraphs if p.text])

def extract_text_from_pptx(file_bytes: bytes) -> str:
    from io import BytesIO
    prs = Presentation(BytesIO(file_bytes))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
    return "\n".join(texts)

def extract_text_from_txt(file_bytes: bytes) -> str:
    try:
        return file_bytes.decode("utf-8")
    except:
        return file_bytes.decode("latin-1", errors="ignore")

def extract_text_from_uploaded_file(uploaded_file) -> Tuple[str, str]:
    name = uploaded_file.name.lower()
    raw = uploaded_file.read()
    if name.endswith(".pdf"):
        return extract_text_from_pdf_bytes(raw), "pdf"
    elif name.endswith(".docx"):
        return extract_text_from_docx(raw), "docx"
    elif name.endswith(".pptx"):
        return extract_text_from_pptx(raw), "pptx"
    elif name.endswith(".txt"):
        return extract_text_from_txt(raw), "txt"
    return "", "unknown"

# -------------------------
# Chunking
# -------------------------
def chunk_text(text: str, chunk_size: int = 2000, overlap: int = 100) -> List[str]:
    text = text.replace("\r", "")
    chunks, start = [], 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        start += chunk_size - overlap
    return chunks

# -------------------------
# Embeddings with caching
# -------------------------
CACHE_DIR = "embedding_cache"
os.makedirs(CACHE_DIR, exist_ok=True)

def get_embeddings_batch(texts: List[str], doc_name: str = None) -> List[List[float]]:
    if doc_name:
        cache_file = os.path.join(CACHE_DIR, doc_name + ".pkl")
        if os.path.exists(cache_file):
            with open(cache_file, "rb") as f:
                return pickle.load(f)

    # Compute embeddings
    embeddings = _embedder.encode(texts, show_progress_bar=True, batch_size=32).tolist()

    # Save cache
    if doc_name:
        with open(os.path.join(CACHE_DIR, doc_name + ".pkl"), "wb") as f:
            pickle.dump(embeddings, f)

    return embeddings

# -------------------------
# Answer Generation (Local or LLM)
# -------------------------
def generate_answer(question: str, contexts: List[dict], model=None) -> str:
    """
    Placeholder: Implement your free model (MPT, WizardLM, etc.) here
    """
    # For now, return concatenated contexts (faster than calling API)
    assembled = "\n\n".join([f"{c['source']}: {c['text']}" for c in contexts])
    return f"Question: {question}\n\nSources:\n{assembled}\n\nAnswer (placeholder)"
